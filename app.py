"""
VERSION 15.2 - CLOUD DEPLOYMENT OPTIMIZED
Enhanced for Streamlit Cloud with better error handling and fallbacks
"""

import streamlit as st
import pandas as pd
from PIL import Image
from io import BytesIO
import openpyxl
from openpyxl.drawing.image import Image as OpenpyxlImage
from pptx import Presentation
from pptx.util import Inches
from bs4 import BeautifulSoup
import time
import re
import requests
from requests.adapters import HTTPAdapter
from urllib3.util.retry import Retry

# --- CONFIGURATION ---
TARGET_WIDTH_PX = 179
TARGET_HEIGHT_PX = 135

def create_robust_session():
    """Create session with aggressive retry and longer timeouts"""
    session = requests.Session()
    
    # More aggressive retry strategy
    retry = Retry(
        total=5,
        backoff_factor=1,
        status_forcelist=[429, 500, 502, 503, 504],
        allowed_methods=["HEAD", "GET", "OPTIONS"]
    )
    
    adapter = HTTPAdapter(max_retries=retry, pool_connections=10, pool_maxsize=10)
    session.mount('http://', adapter)
    session.mount('https://', adapter)
    
    return session

def download_direct_image(image_url, width, height):
    """Download image from direct URL with robust error handling"""
    try:
        print(f"  📥 Downloading: {image_url[:80]}...")
        
        # More comprehensive headers
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
            'Accept': 'image/avif,image/webp,image/apng,image/svg+xml,image/*,*/*;q=0.8',
            'Accept-Language': 'en-US,en;q=0.9',
            'Accept-Encoding': 'gzip, deflate, br',
            'DNT': '1',
            'Connection': 'keep-alive',
            'Upgrade-Insecure-Requests': '1',
        }
        
        session = create_robust_session()
        
        # Make absolute URL if relative
        if image_url.startswith('//'):
            image_url = 'https:' + image_url
        
        r = session.get(image_url, headers=headers, timeout=30)
        r.raise_for_status()
        
        # Verify we got actual image data
        content_type = r.headers.get('content-type', '').lower()
        if 'image' not in content_type and len(r.content) < 1000:
            print(f"  ⚠️ Response doesn't look like an image (Content-Type: {content_type})")
            return None
        
        img = Image.open(BytesIO(r.content))
        
        if img.mode in ("RGBA", "P", "LA"):
            img = img.convert("RGB")
        
        img = img.resize((width, height), Image.Resampling.LANCZOS)
        
        print(f"  ✅ Success ({img.width}x{img.height})")
        return img
        
    except requests.exceptions.Timeout:
        print(f"  ⏱️ Timeout after 30s")
        return None
    except requests.exceptions.RequestException as e:
        print(f"  ❌ Request failed: {e}")
        return None
    except Exception as e:
        print(f"  ❌ Error: {e}")
        return None

def try_scrape_page(url, verbose=True):
    """Try to scrape image from product page with enhanced compatibility"""
    try:
        # Don't try for heavily protected sites
        if "etsy.com" in url or "next.co.uk" in url:
            if verbose:
                print(f"  ⚠️ {url.split('/')[2]} requires manual IMAGE URL")
            return None
        
        if verbose:
            print(f"  🔍 Scraping: {url[:70]}...")
        
        session = create_robust_session()
        
        # Enhanced headers to appear more like real browser
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/120.0.0.0 Safari/537.36',
            'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/avif,image/webp,image/apng,*/*;q=0.8',
            'Accept-Language': 'en-US,en;q=0.9',
            'Accept-Encoding': 'gzip, deflate, br',
            'DNT': '1',
            'Connection': 'keep-alive',
            'Upgrade-Insecure-Requests': '1',
            'Sec-Fetch-Dest': 'document',
            'Sec-Fetch-Mode': 'navigate',
            'Sec-Fetch-Site': 'none',
            'Cache-Control': 'max-age=0',
        }
        
        r = session.get(url, headers=headers, timeout=25, allow_redirects=True)
        
        if verbose:
            print(f"  Status: {r.status_code}, Content-Length: {len(r.content)}")
        
        if r.status_code != 200:
            print(f"  ❌ Bad status: {r.status_code}")
            return None
        
        soup = BeautifulSoup(r.content, 'html.parser')
        
        # Try multiple meta tag strategies
        meta_properties = [
            'og:image',
            'twitter:image',
            'og:image:secure_url',
            'og:image:url',
            'twitter:image:src'
        ]
        
        for prop in meta_properties:
            # Try property attribute
            meta = soup.find("meta", property=prop)
            if not meta:
                # Try name attribute
                meta = soup.find("meta", attrs={"name": prop})
            
            if meta and meta.get("content"):
                img_url = meta["content"]
                
                # Make absolute URL
                if img_url.startswith('//'):
                    img_url = 'https:' + img_url
                elif img_url.startswith('/'):
                    from urllib.parse import urljoin
                    img_url = urljoin(url, img_url)
                
                if verbose:
                    print(f"  ✓ Found via {prop}: {img_url[:60]}...")
                return img_url
        
        # Fallback: Look for JSON-LD product schema
        if verbose:
            print(f"  Trying JSON-LD schema...")
        
        for script in soup.find_all('script', type='application/ld+json'):
            try:
                import json
                data = json.loads(script.string)
                
                # Handle single object or array
                if isinstance(data, list):
                    for item in data:
                        if isinstance(item, dict) and 'image' in item:
                            img = item['image']
                            if isinstance(img, str):
                                if verbose:
                                    print(f"  ✓ Found in JSON-LD: {img[:60]}...")
                                return img
                            elif isinstance(img, list) and len(img) > 0:
                                if verbose:
                                    print(f"  ✓ Found in JSON-LD: {img[0][:60]}...")
                                return img[0]
                elif isinstance(data, dict) and 'image' in data:
                    img = data['image']
                    if isinstance(img, str):
                        if verbose:
                            print(f"  ✓ Found in JSON-LD: {img[:60]}...")
                        return img
                    elif isinstance(img, list) and len(img) > 0:
                        if verbose:
                            print(f"  ✓ Found in JSON-LD: {img[0][:60]}...")
                        return img[0]
            except:
                continue
        
        if verbose:
            print(f"  ❌ No image found in meta tags or JSON-LD")
        
    except requests.exceptions.Timeout:
        print(f"  ⏱️ Timeout after 25s")
    except requests.exceptions.RequestException as e:
        print(f"  ❌ Request error: {type(e).__name__}: {e}")
    except Exception as e:
        print(f"  ❌ Scraping error: {type(e).__name__}: {e}")
    
    return None

def process_row(row, url_col, image_url_col, verbose=True):
    """Process a single row - intelligently handles both columns"""
    
    # Get values from both columns
    url_value = row.get(url_col) if url_col else None
    image_url_value = row.get(image_url_col) if image_url_col else None
    
    # Clean values
    url_value = str(url_value).strip() if pd.notna(url_value) else None
    image_url_value = str(image_url_value).strip() if pd.notna(image_url_value) else None
    
    # Remove 'nan' strings
    if url_value == 'nan':
        url_value = None
    if image_url_value == 'nan':
        image_url_value = None
    
    final_image_url = None
    
    if verbose:
        print(f"\n{'='*70}")
    
    # STRATEGY 1: Check IMAGE URL column first (highest priority)
    if image_url_value:
        if verbose:
            print(f"IMAGE URL: {image_url_value[:70]}...")
        
        # Verify it looks like an image URL
        is_image_url = any(x in image_url_value.lower() for x in 
                          ['.jpg', '.jpeg', '.png', '.gif', '.webp', 'etsystatic', 'xcdn', 'potterybarn', 'cloudfront'])
        
        if is_image_url:
            final_image_url = image_url_value
            if verbose:
                print(f"  ✓ Using IMAGE URL")
        else:
            if verbose:
                print(f"  ⚠️ IMAGE URL doesn't look like image, will try URL column")
    
    # STRATEGY 2: If no IMAGE URL, try URL column
    if not final_image_url and url_value:
        if verbose:
            print(f"URL: {url_value[:70]}...")
        
        # Check if URL itself is a direct image
        if url_value.lower().endswith(('.jpg', '.jpeg', '.png', '.gif', '.webp')):
            final_image_url = url_value
            if verbose:
                print(f"  ✓ URL is direct image")
        else:
            # Try scraping
            final_image_url = try_scrape_page(url_value, verbose)
    
    return final_image_url

# --- STREAMLIT UI ---
st.set_page_config(page_title="Image Automator v15.2", layout="wide")
st.title("📊 Excel & PPT Image Automator v15.2")

st.markdown("""
## 🎯 Cloud-Optimized Version

**Enhanced for Streamlit Cloud deployment:**
- ✅ Better error handling and retries
- ✅ More robust headers for scraping
- ✅ Detailed logging to diagnose issues
- ✅ Longer timeouts for slow connections

**Column Support:**
- **URL**: Product page URLs (auto-scrapes when possible)
- **IMAGE URL**: Direct image URLs (100% reliable)

**Sites that work with auto-scraping:**
- ✅ Pottery Barn, West Elm, Williams Sonoma
- ✅ Target, Walmart
- ✅ Most furniture/home goods sites

**Sites requiring manual IMAGE URL:**
- ❌ Etsy
- ❌ Next.co.uk
---
""")

uploaded_file = st.file_uploader("Upload Excel Template", type=['xlsx'])

if uploaded_file:
    df = pd.read_excel(uploaded_file)
    st.write("### Data Preview")
    st.dataframe(df.head())
    
    # Auto-detect columns
    columns = df.columns.tolist()
    
    url_col_default = None
    image_url_col_default = None
    
    for col in columns:
        col_lower = col.lower()
        if 'image' in col_lower and 'url' in col_lower:
            image_url_col_default = col
        elif col_lower == 'url':
            url_col_default = col
    
    col1, col2, col3 = st.columns(3)
    
    with col1:
        url_col = st.selectbox(
            "URL Column (product pages)", 
            ['(None)'] + columns,
            index=columns.index(url_col_default) + 1 if url_col_default in columns else 0,
            help="Column with product page URLs (auto-scraping)"
        )
        if url_col == '(None)':
            url_col = None
    
    with col2:
        image_url_col = st.selectbox(
            "IMAGE URL Column (direct images)",
            ['(None)'] + columns,
            index=columns.index(image_url_col_default) + 1 if image_url_col_default in columns else 0,
            help="Column with direct image URLs (manual)"
        )
        if image_url_col == '(None)':
            image_url_col = None
    
    with col3:
        target_col_letter = st.text_input("Output Column", value="B").upper()
    
    if not url_col and not image_url_col:
        st.error("⚠️ Select at least one source column")
    else:
        st.success(f"✅ Using: {url_col or 'None'} (URL) + {image_url_col or 'None'} (IMAGE URL)")
    
    # Add verbose logging toggle
    verbose_mode = st.checkbox("Show detailed logs", value=True, 
                               help="Display detailed scraping logs in console")
    
    if st.button("🔍 Test First 3 URLs"):
        if not url_col and not image_url_col:
            st.error("Please select at least one source column")
        else:
            st.write("### Testing...")
            
            success_count = 0
            for i, row in df.head(3).iterrows():
                with st.expander(f"Row {i+1}", expanded=True):
                    final_image_url = process_row(row, url_col, image_url_col, verbose_mode)
                    
                    if final_image_url:
                        img = download_direct_image(final_image_url, TARGET_WIDTH_PX, TARGET_HEIGHT_PX)
                        
                        if img:
                            success_count += 1
                            st.success(f"✅ Success")
                            st.image(img, width=200)
                            st.caption(f"Image URL: {final_image_url[:100]}...")
                        else:
                            st.error(f"❌ Download failed")
                            st.caption(f"Tried: {final_image_url[:100]}...")
                    else:
                        st.error(f"❌ No image URL found")
            
            st.info(f"**Results: {success_count}/3 successful**")
    
    if st.button("Generate Excel"):
        if not url_col and not image_url_col:
            st.error("Please select at least one source column")
        else:
            progress = st.progress(0)
            status = st.empty()
            
            uploaded_file.seek(0)
            wb = openpyxl.load_workbook(uploaded_file)
            ws = wb.active
            
            count = 0
            failed = []
            total = len(df)
            
            for i, row in df.iterrows():
                progress.progress((i + 1) / total)
                
                status.text(f"Processing {i+1}/{total}...")
                
                final_image_url = process_row(row, url_col, image_url_col, verbose=False)
                
                if final_image_url:
                    img = download_direct_image(final_image_url, TARGET_WIDTH_PX, TARGET_HEIGHT_PX)
                    
                    if img:
                        count += 1
                        buf = BytesIO()
                        img.save(buf, format='PNG')
                        buf.seek(0)
                        excel_img = OpenpyxlImage(buf)
                        ws.add_image(excel_img, f"{target_col_letter}{i+2}")
                        ws.row_dimensions[i+2].height = 105
                    else:
                        failed.append(i+2)
                else:
                    failed.append(i+2)
                        
                time.sleep(0.8)  # Slightly longer delay for cloud
            
            if failed:
                status.warning(f"⚠️ {count}/{total} images added. Failed: {len(failed)} rows")
            else:
                status.success(f"✅ Complete! {count}/{total} images added")
            
            out = BytesIO()
            wb.save(out)
            out.seek(0)
            
            st.download_button("📥 Download Excel", out, "output_v15.2.xlsx",
                              mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")
            
            if failed and len(failed) <= 20:
                st.info(f"**Failed rows:** {', '.join(map(str, failed))}")
    
    if st.button("Generate PowerPoint"):
        if not url_col and not image_url_col:
            st.error("Please select at least one source column")
        else:
            with st.spinner("Generating slides..."):
                prs = Presentation()
                blank = prs.slide_layouts[6]
                count = 0
                
                for i, row in df.iterrows():
                    final_image_url = process_row(row, url_col, image_url_col, verbose=False)
                    
                    if final_image_url:
                        img = download_direct_image(final_image_url, 800, 600)
                        
                        if img:
                            count += 1
                            slide = prs.slides.add_slide(blank)
                            buf = BytesIO()
                            img.save(buf, format='PNG')
                            buf.seek(0)
                            slide.shapes.add_picture(buf, Inches(1), Inches(1), width=Inches(8))
                            
                            tx = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
                            tx.text_frame.text = f"Source: {final_image_url[:100]}"
                
                out = BytesIO()
                prs.save(out)
                out.seek(0)
                st.download_button("📥 Download PowerPoint", out, "output_v15.2.pptx")
                st.success(f"✅ Created {count} slides")
